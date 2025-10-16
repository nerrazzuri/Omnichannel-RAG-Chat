"""
Document processing: chunking and embedding using OpenAI.
"""
from typing import List, Tuple, Dict, Any
from sqlalchemy.orm import Session
from shared.database.models import Document, KnowledgeChunk, KnowledgeBase, Tenant
from shared.utils.storage import write_metadata
from openai import OpenAI
import os, hashlib, struct, random
import io
import csv
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from shared.vector.qdrant import qdrant_service
import logging
import re


def chunk_text(text: str, chunk_size: int = 700, overlap: int = 100) -> List[str]:
    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = max(0, end - overlap)
    return chunks


class DocumentService:
    def __init__(self, db: Session):
        self.db = db
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    @staticmethod
    def _split_sentences(text: str) -> List[str]:
        # Lightweight sentence splitter using punctuation and line breaks
        # Avoid breaking on common abbreviations by a simple heuristic
        text = re.sub(r"\s+", " ", text)
        candidates = re.split(r"(?<=[\.!?])\s+(?=[A-Z(\[])", text)
        sentences: List[str] = []
        for s in candidates:
            s = s.strip()
            if not s:
                continue
            if len(s) < 3:
                continue
            sentences.append(s)
        return sentences

    def _build_chunks_with_metadata(self, text: str, target_chars: int = 1400, overlap_sentences: int = 2) -> List[Tuple[str, Dict[str, Any]]]:
        """Sentence-aware chunking with small overlap and chapter/page tagging.

        Recognizes page markers like [[PAGE:n]] if present.
        Detects chapter headings like "Chapter 3. ..." and tags subsequent chunks
        until the next heading.
        """
        # Detect simple page markers
        pages: List[Tuple[int, str]] = []
        page_matches = list(re.finditer(r"\[\[PAGE:(\d+)\]\]", text))
        if page_matches:
            last_idx = 0
            current_page = 1
            segments: List[Tuple[int, str]] = []
            for m in page_matches:
                seg = text[last_idx:m.start()]
                if seg.strip():
                    segments.append((current_page, seg))
                try:
                    current_page = int(m.group(1))
                except Exception:
                    pass
                last_idx = m.end()
            tail = text[last_idx:]
            if tail.strip():
                segments.append((current_page, tail))
            pages = segments
        else:
            pages = [(None, text)]  # type: ignore

        chunks: List[Tuple[str, Dict[str, Any]]] = []
        current_chapter_num: Any = None
        current_chapter_title: str = ""

        for page_num, page_text in pages:
            # Identify chapter heading at start of page or within first lines
            for line in page_text.splitlines()[:6]:
                m = re.match(r"^\s*chapter\s+(\d+)\s*[\.:\-]?\s*(.*)$", line.strip(), flags=re.IGNORECASE)
                if m:
                    try:
                        current_chapter_num = int(m.group(1))
                        current_chapter_title = (m.group(2) or "").strip()
                    except Exception:
                        pass
                    break

            sentences = self._split_sentences(page_text)
            if not sentences:
                continue

            buf: List[str] = []
            for i, s in enumerate(sentences):
                if not buf:
                    buf.append(s)
                else:
                    prospective = (" ".join(buf) + " " + s).strip()
                    if len(prospective) <= target_chars:
                        buf.append(s)
                    else:
                        # Emit chunk
                        text_chunk = " ".join(buf).strip()
                        meta: Dict[str, Any] = {}
                        if page_num is not None:
                            meta["page"] = page_num
                        if current_chapter_num is not None:
                            meta["chapter_num"] = current_chapter_num
                        if current_chapter_title:
                            meta["chapter_title"] = current_chapter_title
                        chunks.append((text_chunk, meta))
                        # Start new buffer with overlap
                        overlap = sentences[max(0, i - overlap_sentences):i]
                        buf = overlap + [s]

            if buf:
                text_chunk = " ".join(buf).strip()
                meta: Dict[str, Any] = {}
                if page_num is not None:
                    meta["page"] = page_num
                if current_chapter_num is not None:
                    meta["chapter_num"] = current_chapter_num
                if current_chapter_title:
                    meta["chapter_title"] = current_chapter_title
                chunks.append((text_chunk, meta))

        return chunks

    @staticmethod
    def _extract_chapter_info(text: str) -> Dict[str, Any]:
        """Extract simple chapter metadata from a text chunk.

        Looks for patterns like:
        - "Chapter 3. Program and Project Management Roles and Responsibilities"
        - "Chapter 4: Governance"
        """
        try:
            import re
            lines = [l.strip() for l in text.splitlines() if l.strip()]
            for line in lines[:5]:  # inspect only early lines of the chunk
                m = re.match(r"^chapter\s+(\d+)\s*[\.:\-]?\s*(.*)$", line, flags=re.IGNORECASE)
                if m:
                    num = int(m.group(1))
                    title = (m.group(2) or "").strip()
                    return {"chapter_num": num, "chapter_title": title}
        except Exception:
            pass
        return {}

    def embed(self, inputs: List[str]) -> List[List[float]]:
        api_key = os.getenv("OPENAI_API_KEY")
        if api_key:
            # Batch requests to respect OpenAI per-request token limits
            def estimate_tokens(text: str) -> int:
                # Rough heuristic: 4 chars per token
                return max(1, len(text) // 4)

            MAX_TOKENS_PER_REQUEST = 280_000  # keep below 300k limit
            embeddings: List[List[float]] = []
            batch: List[str] = []
            tokens_in_batch = 0
            for t in inputs:
                t_tokens = estimate_tokens(t)
                if batch and tokens_in_batch + t_tokens > MAX_TOKENS_PER_REQUEST:
                    resp = self.client.embeddings.create(
                        model="text-embedding-3-small",
                        input=batch,
                    )
                    embeddings.extend([d.embedding for d in resp.data])
                    batch = []
                    tokens_in_batch = 0
                batch.append(t)
                tokens_in_batch += t_tokens

            if batch:
                resp = self.client.embeddings.create(
                    model="text-embedding-3-small",
                    input=batch,
                )
                embeddings.extend([d.embedding for d in resp.data])

            return embeddings
        # Fallback deterministic embedding (no external dependency)
        vectors: List[List[float]] = []
        dim = 256
        for text in inputs:
            h = hashlib.sha256(text.encode("utf-8")).digest()
            # Expand hash deterministically
            rnd = random.Random(h)
            vec = [rnd.uniform(-1.0, 1.0) for _ in range(dim)]
            vectors.append(vec)
        return vectors

    def process_and_store(self, tenant_id: str, title: str, content: str, knowledge_base_id: str) -> Tuple[str, int]:
        try:
            # Validate tenant_id is a valid UUID
            import uuid
            try:
                uuid.UUID(tenant_id)
            except ValueError:
                raise ValueError(f"Invalid tenant_id: {tenant_id}. Must be a valid UUID.")
            
            # Ensure a knowledge base exists for this tenant
            kb_id = self._get_or_create_knowledge_base(tenant_id, knowledge_base_id)

            # Create document
            doc = Document(title=title, content=content, knowledge_base_id=kb_id, status="PROCESSING")
            self.db.add(doc)
            self.db.commit()
            self.db.refresh(doc)

            # Chunk and embed (sentence-aware, chapter-aware)
            chunk_pairs = self._build_chunks_with_metadata(content)
            chunks = [t for (t, _m) in chunk_pairs]
            metas = [m for (_t, m) in chunk_pairs]
            if not chunks:
                raise ValueError("No chunks could be created from the content")
            
            embeddings = self.embed(chunks)

            # Store chunks
            qdrant_payload: List[Dict[str, Any]] = []
            for idx, ((chunk_text_val, meta_chunk), emb) in enumerate(zip(chunk_pairs, embeddings)):
                # Ensure a concrete UUID is assigned before using the ID
                import uuid as _uuid
                chunk_id = _uuid.uuid4()
                # Merge auto-headline detection with chunk-derived meta
                chapter_meta = self._extract_chapter_info(chunk_text_val)
                merged_meta = dict(meta_chunk)
                for k, v in chapter_meta.items():
                    if v and k not in merged_meta:
                        merged_meta[k] = v
                kc = KnowledgeChunk(
                    id=chunk_id,
                    document_id=doc.id,
                    content=chunk_text_val,
                    chunk_index=idx,
                    embedding=emb,
                    meta=merged_meta or {},
                )
                self.db.add(kc)
                # SQLAlchemy default UUID is assigned on instantiation; id is available before commit
                try:
                    qdrant_payload.append({
                        "id": str(chunk_id),
                        "embedding": emb,
                        "document_id": str(doc.id),
                        "content": chunk_text_val,
                        "chunk_index": idx,
                        "chapter_num": merged_meta.get("chapter_num"),
                        "chapter_title": merged_meta.get("chapter_title"),
                        "page": merged_meta.get("page"),
                    })
                except Exception:
                    pass
            doc.status = "INDEXED"
            doc.chunk_count = len(chunks)
            self.db.add(doc)
            self.db.commit()

            # Best-effort: upsert vectors to Qdrant when dimensions match (OpenAI = 1536)
            try:
                if qdrant_payload and isinstance(qdrant_payload[0].get("embedding"), list):
                    dim = len(qdrant_payload[0]["embedding"]) if qdrant_payload[0].get("embedding") else 0
                    if dim == 1536:
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            # Collection may already exist or service may be unavailable
                            pass
                        try:
                            qdrant_service.upsert_knowledge_chunks(tenant_id, qdrant_payload)
                        except Exception as e:
                            logging.getLogger(__name__).warning(f"Qdrant upsert skipped: {e}")
                    else:
                        logging.getLogger(__name__).info("Skipping Qdrant upsert due to embedding dimension mismatch")
            except Exception:
                # Never fail ingestion due to vector store issues
                pass
            
            # Write metadata.json for downstream processing
            metadata: Dict[str, Any] = {
                "tenant_id": tenant_id,
                "document_id": str(doc.id),
                "knowledge_base_id": kb_id,
                "title": title,
                "chunk_count": doc.chunk_count,
                "status": doc.status,
            }
            base_path = os.getenv("DOCUMENT_STORAGE_PATH", os.path.join(os.getcwd(), "storage"))
            try:
                write_metadata(base_path, tenant_id, str(doc.id), metadata)
            except Exception as e:
                # Log but don't fail if metadata write fails
                import logging
                logger = logging.getLogger(__name__)
                logger.warning(f"Failed to write metadata: {e}")
            
            return str(doc.id), len(chunks)
        except Exception as e:
            # If there's an error, rollback the transaction
            self.db.rollback()
            raise

    def _get_or_create_knowledge_base(self, tenant_id: str, provided_kb_id: str) -> str:
        import uuid
        
        # Validate tenant_id
        try:
            tenant_uuid = uuid.UUID(tenant_id)
        except ValueError:
            raise ValueError(f"Invalid tenant_id format: {tenant_id}")
        
        # Ensure tenant exists; if missing, create a default record (dev-friendly)
        try:
            existing_tenant = self.db.get(Tenant, tenant_uuid)
            if not existing_tenant:
                # Create a minimal tenant to satisfy FK; name/domain deterministic for the given UUID
                t = Tenant(id=tenant_uuid, name="Seeded Tenant", domain="seeded", settings={})
                self.db.add(t)
                self.db.commit()
        except Exception:
            # Best-effort; if this fails, the subsequent KB creation will surface the error
            self.db.rollback()
        
        # If a valid kb id is provided and exists, use it
        if provided_kb_id and provided_kb_id != "00000000-0000-0000-0000-000000000000":
            try:
                kb_uuid = uuid.UUID(provided_kb_id)
                kb = self.db.get(KnowledgeBase, kb_uuid)
                if kb:
                    return str(kb.id)
            except ValueError:
                # Invalid UUID format, skip
                pass

        # Try to find default KB for tenant
        kb = (
            self.db.query(KnowledgeBase)
            .filter(KnowledgeBase.tenant_id == tenant_uuid)
            .order_by(KnowledgeBase.created_at.asc())
            .first()
        )
        if kb:
            return str(kb.id)

        # Create a new knowledge base for this tenant
        new_kb = KnowledgeBase(tenant_id=tenant_uuid, name="Default", status="ACTIVE", document_count=0)
        self.db.add(new_kb)
        self.db.commit()
        self.db.refresh(new_kb)
        return str(new_kb.id)

    def extract_text_from_file(self, filename: str, data: bytes) -> str:
        name = filename.lower()
        if name.endswith('.txt') or name.endswith('.csv'):
            try:
                return data.decode('utf-8')
            except Exception:
                return data.decode('latin-1', errors='ignore')
        if name.endswith('.docx'):
            buf = io.BytesIO(data)
            d = DocxDocument(buf)
            return '\n'.join(p.text for p in d.paragraphs)
        if name.endswith('.pptx'):
            buf = io.BytesIO(data)
            prs = Presentation(buf)
            texts: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        texts.append(shape.text)
            return '\n'.join(texts)
        if name.endswith('.xlsx'):
            buf = io.BytesIO(data)
            wb = load_workbook(buf, data_only=True)
            texts: List[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    texts.append('\t'.join('' if v is None else str(v) for v in row))
            return '\n'.join(texts)
        # PDF handled by PyPDF2
        if name.endswith('.pdf'):
            from PyPDF2 import PdfReader
            buf = io.BytesIO(data)
            reader = PdfReader(buf)
            texts: List[str] = []
            for i, page in enumerate(reader.pages, start=1):
                texts.append(f"[[PAGE:{i}]]\n" + (page.extract_text() or ''))
            return '\n'.join(texts)
        # Fallback raw decode
        try:
            return data.decode('utf-8')
        except Exception:
            return data.decode('latin-1', errors='ignore')

    def extract_rows_from_file(self, filename: str, data: bytes) -> List[str]:
        name = filename.lower()
        rows: List[str] = []
        if name.endswith('.csv'):
            try:
                text = data.decode('utf-8-sig', errors='ignore')
            except Exception:
                text = data.decode('latin-1', errors='ignore')
            reader = csv.reader(io.StringIO(text))
            # Re-serialize each row via csv.writer to preserve quoting and commas inside fields
            for r in reader:
                buf = io.StringIO()
                writer = csv.writer(buf)
                writer.writerow(r)
                rows.append(buf.getvalue().strip('\n'))
        elif name.endswith('.xlsx'):
            buf = io.BytesIO(data)
            wb = load_workbook(buf, data_only=True)
            for ws in wb.worksheets:
                for r in ws.iter_rows(values_only=True):
                    # Use csv.writer to serialize each row consistently
                    buf_row = io.StringIO()
                    writer = csv.writer(buf_row)
                    writer.writerow(['' if v is None else v for v in r])
                    rows.append(buf_row.getvalue().strip('\n'))
        return rows

    def process_rows_and_store(self, tenant_id: str, title: str, rows: List[str], knowledge_base_id: str) -> Tuple[str, int]:
        try:
            # Validate tenant_id
            import uuid
            try:
                uuid.UUID(tenant_id)
            except ValueError:
                raise ValueError(f"Invalid tenant_id: {tenant_id}. Must be a valid UUID.")
            
            if not rows:
                raise ValueError("No rows provided to process")
            
            kb_id = self._get_or_create_knowledge_base(tenant_id, knowledge_base_id)
            preview = '\n'.join(rows[:5]) + ('\n...' if len(rows) > 5 else '')
            doc = Document(title=title, content=preview, knowledge_base_id=kb_id, status="PROCESSING")
            
            # Capture header columns if present (first row)
            if rows:
                try:
                    header_reader = csv.reader(io.StringIO(rows[0]))
                    header = next(header_reader)
                    doc.meta = {"columns": [h.strip().lower() for h in header]}
                    # If first row is header, skip it for chunk storage
                    data_rows = rows[1:]
                except Exception:
                    data_rows = rows
            else:
                data_rows = rows
            
            if not data_rows:
                raise ValueError("No data rows found after header extraction")
            
            self.db.add(doc)
            self.db.commit()
            self.db.refresh(doc)

            embeddings = self.embed(data_rows)
            qdrant_payload: List[Dict[str, Any]] = []
            for idx, (row_text, emb) in enumerate(zip(data_rows, embeddings)):
                import uuid as _uuid
                chunk_id = _uuid.uuid4()
                # Tabular rows do not carry chapter info
                kc = KnowledgeChunk(id=chunk_id, document_id=doc.id, content=row_text, chunk_index=idx, embedding=emb)
                self.db.add(kc)
                try:
                    qdrant_payload.append({
                        "id": str(chunk_id),
                        "embedding": emb,
                        "document_id": str(doc.id),
                        "content": row_text,
                        "chunk_index": idx,
                        "chapter_num": None,
                        "chapter_title": None,
                    })
                except Exception:
                    pass
            doc.status = "INDEXED"
            doc.chunk_count = len(data_rows)
            self.db.add(doc)
            self.db.commit()

            # Best-effort: upsert to Qdrant when dimensions match expected size
            try:
                if qdrant_payload and isinstance(qdrant_payload[0].get("embedding"), list):
                    dim = len(qdrant_payload[0]["embedding"]) if qdrant_payload[0].get("embedding") else 0
                    if dim == 1536:
                        try:
                            qdrant_service.create_collection()
                        except Exception:
                            pass
                        try:
                            qdrant_service.upsert_knowledge_chunks(tenant_id, qdrant_payload)
                        except Exception as e:
                            logging.getLogger(__name__).warning(f"Qdrant upsert skipped: {e}")
                    else:
                        logging.getLogger(__name__).info("Skipping Qdrant upsert due to embedding dimension mismatch")
            except Exception:
                pass

            metadata: Dict[str, Any] = {
                "tenant_id": tenant_id,
                "document_id": str(doc.id),
                "knowledge_base_id": kb_id,
                "title": title,
                "chunk_count": doc.chunk_count,
                "status": doc.status,
            }
            base_path = os.getenv("DOCUMENT_STORAGE_PATH", os.path.join(os.getcwd(), "storage"))
            try:
                write_metadata(base_path, tenant_id, str(doc.id), metadata)
            except Exception as e:
                # Log but don't fail if metadata write fails
                import logging
                logger = logging.getLogger(__name__)
                logger.warning(f"Failed to write metadata: {e}")
            
            return str(doc.id), len(rows)
        except Exception as e:
            # Rollback on error
            self.db.rollback()
            raise


